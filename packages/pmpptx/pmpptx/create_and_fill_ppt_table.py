from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_and_fill_ppt_table(row_names, col_names, slide, act_df, yoy_df, yoy2_df, prs, titles = 'none', act_font_size = Pt(16), yoy_font_size = Pt(12), yoy_gray = RGBColor(128, 128, 128)):

    shapes = slide.shapes

    l, t, wi, he = Inches(1), Inches(.25), Inches(8), Inches(2)

    shapes[1].left = l
    shapes[1].top = t
    shapes[1].width = wi
    shapes[1].height = he

    title_idx = prs.slides.index(slide) - 1

    for shape in shapes:
        if shape.placeholder_format.type == 1:
            if titles != 'none':
                shape.text = titles[title_idx]

    x, y, cx, cy = Inches(1), Inches(.75), Inches(8), Inches(5)

    nrow = len(row_names) + 1
    ncol = len(col_names)

    shape = shapes.add_table(nrow, ncol, x, y, cx, cy)

    table = shape.table

    for i in range(0, nrow - 1):
        update_cell = table.cell(i + 1, 0)
        update_cell.text_frame.paragraphs[0].text = row_names[i]
        update_cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        update_cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


    for j in range(0, ncol):
        update_cell = table.cell(0, j)
        update_cell.text_frame.paragraphs[0].text = col_names[j]
        update_cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        update_cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    for i in range(0, nrow - 1):
        for j in range(0, ncol - 2):
            update_cell = table.cell(i + 1, j + 1)
            up = update_cell.text_frame.paragraphs[0]
            r_num = up.add_run()
            r_num.text = act_df.iloc[j, (i + 1)]
            r_num.font.size = act_font_size
            r_per = up.add_run()
            r_per.text = yoy_df.iloc[j, (i + 1)]
            r_per.font.color.rgb = yoy_gray
            r_per.font.size = yoy_font_size
            update_cell.vertical_anchor = MSO_ANCHOR.TOP
            update_cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    for i in range(0, nrow - 1):
        update_cell = table.cell(i + 1, ncol - 1)
        update_cell.text = yoy2_df.iloc[0, i + 1]
        update_cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        update_cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
